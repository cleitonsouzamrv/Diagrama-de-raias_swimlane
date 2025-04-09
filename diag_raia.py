import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from io import BytesIO
import textwrap
import zipfile
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

# Configuração da página
st.set_page_config(page_title="Swimlane Generator", layout="wide")
st.title("📊 Gerador de Diagrama de Raias (Swimlane)")

@st.cache_resource(show_spinner="Carregando planilha...")
def load_excel(file):
    return pd.ExcelFile(file)

@st.cache_data(show_spinner="Lendo guia...")
def load_sheet(file, sheet_name):
    return pd.read_excel(file, sheet_name=sheet_name)

def ajustar_fonte_em_caixa_fixa(texto, largura_max_chars=28, box_height=0.8, font_min=4, font_max=10):
    wrapped = textwrap.wrap(texto, width=largura_max_chars)
    num_linhas = len(wrapped)
    line_spacing = box_height / (num_linhas + 1)
    font_size = min(max(int(line_spacing * 10), font_min), font_max)
    return wrapped, font_size, line_spacing

def gerar_diagrama(df_filtrado, col_funcao, col_atividade, col_ordem, valor_filtro, nome_titulo, col_cor=None):
    cols = [col_funcao, col_atividade, col_ordem]
    if col_cor:
        cols.append(col_cor)
    df_filtrado = df_filtrado[cols].dropna()
    df_filtrado[col_ordem] = df_filtrado[col_ordem].astype(str)
    funcoes = df_filtrado[col_funcao].unique().tolist()
    sequencias = sorted(df_filtrado[col_ordem].unique().tolist())
    fun_y = {f: i for i, f in enumerate(reversed(funcoes))}
    seq_x = {s: i for i, s in enumerate(sequencias)}
    width = max(12, len(seq_x) * 1.5)
    height = max(6, len(fun_y) * 0.9)
    fig, ax = plt.subplots(figsize=(width, height), facecolor='white')
    for func, y in fun_y.items():
        ax.add_patch(patches.Rectangle((0, y), len(seq_x), 1, color="#F5F5F5", ec="black", linewidth=0.8))
        ax.text(-0.3, y + 0.5, func, va="center", ha="right", fontsize=10, fontweight="bold", color="#333333")
    for _, row in df_filtrado.iterrows():
        x = seq_x[row[col_ordem]]
        y = fun_y[row[col_funcao]]
        desc = str(row[col_atividade])
        cor_caixa = row[col_cor] if col_cor and pd.notna(row[col_cor]) else "#034E2B"
        box_width = 0.9
        box_height = 0.8
        wrapped, font_size, line_spacing = ajustar_fonte_em_caixa_fixa(desc, largura_max_chars=28, box_height=box_height)
        ax.add_patch(patches.Rectangle((x + 0.05, y + 0.1), box_width, box_height, color=cor_caixa, ec=None))
        for i, line in enumerate(wrapped):
            y_pos = y + 0.1 + box_height - ((i + 1) * line_spacing)
            ax.text(x + 0.5, y_pos, line, ha="center", va="center", fontsize=font_size, color="white")
    ax.set_xlim(-1, len(seq_x))
    ax.set_ylim(0, len(fun_y))
    ax.set_xticks(range(len(seq_x)))
    ax.set_xticklabels(sequencias, rotation=45, ha='right', fontsize=7)
    ax.set_yticks([])
    ax.set_title(f"Diagrama Swimlane – {nome_titulo} = {valor_filtro}", fontsize=13, color="#333333")
    ax.axis("off")
    buffer = BytesIO()
    plt.tight_layout()
    fig.savefig(buffer, format="png", dpi=300, bbox_inches='tight', facecolor='white')
    buffer.seek(0)
    return buffer

uploaded_file = st.file_uploader("📎 Faça upload de uma planilha Excel (.xlsx)", type=["xlsx"])

if uploaded_file:
    excel = load_excel(uploaded_file)
    sheet_name = st.selectbox("🗂️ Escolha a guia da planilha", excel.sheet_names)
    df = load_sheet(uploaded_file, sheet_name)

    st.markdown("#### 🔍 Pré-visualização")
    st.dataframe(df.head(6))

    col_funcao = st.selectbox("🧑‍💼 Coluna da Raia (FUNÇÃO)", df.columns)
    col_atividade = st.selectbox("📝 Coluna da Descrição da Atividade (DESCRIÇÃO DO PROCESSO)", df.columns)
    col_ordem = st.selectbox("🔢 Coluna da Ordem (SEQUÊNCIA)", df.columns)
    col_filtro = st.selectbox("Coluna para Diagrama (ATIVIDADE (FLUXO))", df.columns)

    col_cor = st.selectbox("🎨 Coluna com as cores (hex)", [None] + list(df.columns), help="Opcional: selecione a coluna com códigos hex para colorir as atividades")

    valores_filtro = sorted(df[col_filtro].dropna().astype(str).unique().tolist())
    valor_escolhido = st.selectbox("Escolha a Descrição da Atividade (ATIVIDADE PARA GERAR O DIAGRAMA)", valores_filtro)

    col1, col2, col3 = st.columns(3)

    if col1.button("🚀 Gerar Diagrama Selecionado"):
        df_filtrado = df[df[col_filtro].astype(str) == valor_escolhido]
        if df_filtrado.empty:
            st.warning("⚠️ Nenhuma atividade encontrada para esse filtro.")
            st.stop()
        buffer = gerar_diagrama(df_filtrado, col_funcao, col_atividade, col_ordem, valor_escolhido, col_filtro, col_cor)
        st.markdown("### 🖼️ Diagrama de Raia Gerado")
        st.image(buffer, caption="Visualização do Diagrama Swimlane", use_container_width=True)
        st.download_button("📥 Baixar imagem do Diagrama (PNG)", data=buffer, file_name=f"diagrama_swimlane_{valor_escolhido}.png", mime="image/png")

    if col2.button("📊 Gerar todos os Diagramas (.pptx)"):
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        for valor in valores_filtro:
            df_filtrado = df[df[col_filtro].astype(str) == valor]
            if df_filtrado.empty:
                continue
            buffer = gerar_diagrama(df_filtrado, col_funcao, col_atividade, col_ordem, valor, col_filtro, col_cor)
            slide = prs.slides.add_slide(blank_slide_layout)
            safe_valor = "".join(c if c.isalnum() or c in " _-" else "_" for c in valor)
            title_shape = slide.shapes.title if slide.shapes.title else None
            if not title_shape:
                title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
                title_shape.text = valor
                title_shape.text_frame.paragraphs[0].font.size = Pt(20)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            img_path = f"temp_{safe_valor}.png"
            with open(img_path, "wb") as f:
                f.write(buffer.getvalue())
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(0.9), width=Inches(9))
            os.remove(img_path)
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        st.download_button("📥 Baixar todos os Diagramas em PPTX", data=pptx_buffer, file_name="diagramas_swimlane.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    if col2.button("🖼️ Gerar todos os Diagramas em PNG (.zip)"):
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w") as zipf:
            for valor in valores_filtro:
                df_filtrado = df[df[col_filtro].astype(str) == valor]
                if df_filtrado.empty:
                    continue
                buffer = gerar_diagrama(df_filtrado, col_funcao, col_atividade, col_ordem, valor, col_filtro, col_cor)
                safe_valor = "".join(c if c.isalnum() or c in " _-" else "_" for c in valor)
                zipf.writestr(f"{safe_valor}.png", buffer.getvalue())
        zip_buffer.seek(0)
        st.download_button("📥 Baixar todos os Diagramas (PNG .zip)", data=zip_buffer, file_name="diagramas_swimlane.zip", mime="application/zip")
